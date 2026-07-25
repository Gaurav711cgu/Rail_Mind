import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN

def create_slide(prs):
    # Uses a blank slide layout to draw custom styled layouts
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    
    # Set dark space background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(10, 16, 26)
    return slide

def add_header(slide, title_text, subtitle_text=""):
    # Title Text Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text.upper()
    p.font.name = "Segoe UI"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0, 212, 255) # Neon Cyan
    
    # Subtitle
    if subtitle_text:
        sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(11.7), Inches(0.5))
        tf_sub = sub_box.text_frame
        tf_sub.word_wrap = True
        tf_sub.margin_left = tf_sub.margin_top = tf_sub.margin_right = tf_sub.margin_bottom = 0
        p_sub = tf_sub.paragraphs[0]
        p_sub.text = subtitle_text
        p_sub.font.name = "Segoe UI"
        p_sub.font.size = Pt(14)
        p_sub.font.color.rgb = RGBColor(140, 150, 170) # Muted Text
        
    # Top horizontal divider line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.6), Inches(11.733), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(30, 41, 59) # Dark border color
    line.line.fill.background()

def draw_card(slide, left, top, width, height, fill_color, border_color, border_width_pt=1):
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = border_color
    card.line.width = Pt(border_width_pt)
    try:
        card.adjustments[0] = 0.04  # Subtle round corners
    except Exception:
        pass
    return card

def add_bullet(tf, bold_prefix, text_content, font_size=13, color=RGBColor(220, 225, 235)):
    p = tf.add_paragraph()
    p.space_after = Pt(8)
    p.font.name = "Segoe UI"
    p.font.size = Pt(font_size)
    
    run_bold = p.add_run()
    run_bold.text = bold_prefix
    run_bold.font.bold = True
    run_bold.font.color.rgb = RGBColor(0, 212, 255) # Cyan prefix
    
    run_text = p.add_run()
    run_text.text = text_content
    run_text.font.color.rgb = color

def build_deck():
    print("Initializing Elite PowerPoint Presentation Engine...")
    prs = Presentation()
    
    # Configure 16:9 widescreen slides
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # Brand Color Palette
    COLOR_BG = RGBColor(10, 16, 26)       # Deep Space Black
    COLOR_CARD = RGBColor(18, 25, 41)     # Dark Slate Blue
    COLOR_CYAN = RGBColor(0, 212, 255)     # Neon Cyan Accent
    COLOR_BLUE = RGBColor(30, 144, 255)    # Royal Blue
    COLOR_WHITE = RGBColor(240, 244, 248)  # High Contrast White
    COLOR_MUTED = RGBColor(140, 150, 170)  # Cool Gray
    COLOR_RED = RGBColor(255, 99, 99)      # Delay Warning Red
    COLOR_GREEN = RGBColor(0, 230, 118)    # Audit Success Green

    # ==========================================
    # SLIDE 1: Title Slide
    # ==========================================
    slide_1 = create_slide(prs)
    
    # Glowing Central Box
    draw_card(slide_1, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5), COLOR_CARD, COLOR_CYAN, border_width_pt=2)
    
    # Title Text Box
    title_box = slide_1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = "RAILMIND"
    p.font.name = "Segoe UI"
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    # Subtitle
    sub_box = slide_1.shapes.add_textbox(Inches(1.5), Inches(3.6), Inches(10.333), Inches(1.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.alignment = PP_ALIGN.CENTER
    p_sub.text = "Autonomous Decision Grid & Real-Time Telemetry for Indian Railways"
    p_sub.font.name = "Segoe UI"
    p_sub.font.size = Pt(20)
    p_sub.font.color.rgb = COLOR_WHITE
    
    p_sub2 = tf_sub.add_paragraph()
    p_sub2.alignment = PP_ALIGN.CENTER
    p_sub2.space_before = Pt(8)
    p_sub2.text = "Delhi Regionals Screening | India's Biggest International Hackathon — Far Away 2026"
    p_sub2.font.name = "Segoe UI"
    p_sub2.font.size = Pt(13)
    p_sub2.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 2: The Operational Crisis (Problem)
    # ==========================================
    slide_2 = create_slide(prs)
    add_header(slide_2, "The Operational Crisis: Indian Railways", "Why traditional manual scheduling systems drop efficiency under saturation load")
    
    # Card 1: Punctuality Collapse
    draw_card(slide_2, Inches(0.8), Inches(2.0), Inches(3.64), Inches(4.5), COLOR_CARD, COLOR_RED, border_width_pt=1)
    # Title
    t1 = slide_2.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(3.24), Inches(0.6))
    t1.text_frame.word_wrap = True
    p = t1.text_frame.paragraphs[0]
    p.text = "PUNCTUALITY SHIFT"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    # Value
    v1 = slide_2.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(3.24), Inches(1.0))
    p = v1.text_frame.paragraphs[0]
    p.text = "73.62%"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    # Desc
    d1 = slide_2.shapes.add_textbox(Inches(1.0), Inches(4.0), Inches(3.24), Inches(2.2))
    d1.text_frame.word_wrap = True
    p = d1.text_frame.paragraphs[0]
    p.text = "Corridor congestion and line saturation have led to a precipitous drop in overall punctuality. When one train halts, the delay cascade transmits across adjacent blocks unchecked."
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED
    
    # Card 2: Manual Heuristics
    draw_card(slide_2, Inches(4.84), Inches(2.0), Inches(3.64), Inches(4.5), COLOR_CARD, COLOR_RED, border_width_pt=1)
    # Title
    t2 = slide_2.shapes.add_textbox(Inches(5.04), Inches(2.2), Inches(3.24), Inches(0.6))
    t2.text_frame.word_wrap = True
    p = t2.text_frame.paragraphs[0]
    p.text = "HUMAN HEURISTICS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    # Value
    v2 = slide_2.shapes.add_textbox(Inches(5.04), Inches(2.8), Inches(3.24), Inches(1.0))
    p = v2.text_frame.paragraphs[0]
    p.text = "Manual"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    # Desc
    d2 = slide_2.shapes.add_textbox(Inches(5.04), Inches(4.0), Inches(3.24), Inches(2.2))
    d2.text_frame.word_wrap = True
    p = d2.text_frame.paragraphs[0]
    p.text = "Section controllers currently resolve complex scheduling overlaps and loop holdings using personal bias and mental calculations under extreme stress, creating sub-optimal grid lockouts."
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED

    # Card 3: Controllable Delay Factors
    draw_card(slide_2, Inches(8.88), Inches(2.0), Inches(3.64), Inches(4.5), COLOR_CARD, COLOR_RED, border_width_pt=1)
    # Title
    t3 = slide_2.shapes.add_textbox(Inches(9.08), Inches(2.2), Inches(3.24), Inches(0.6))
    t3.text_frame.word_wrap = True
    p = t3.text_frame.paragraphs[0]
    p.text = "CONTROLLABLE LOSS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_RED
    # Value
    v3 = slide_2.shapes.add_textbox(Inches(9.08), Inches(2.8), Inches(3.24), Inches(1.0))
    p = v3.text_frame.paragraphs[0]
    p.text = "27 of 33"
    p.font.name = "Segoe UI"
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = COLOR_WHITE
    # Desc
    d3 = slide_2.shapes.add_textbox(Inches(9.08), Inches(4.0), Inches(3.24), Inches(2.2))
    d3.text_frame.word_wrap = True
    p = d3.text_frame.paragraphs[0]
    p.text = "A majority of delay-causing factors on corridors are predictable and controllable, yet they lack real-time digital automation feedback loops to dynamically reroute and optimize dispatching."
    p.font.name = "Segoe UI"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_MUTED

    # ==========================================
    # SLIDE 3: RailMind Solution
    # ==========================================
    slide_3 = create_slide(prs)
    add_header(slide_3, "RailMind: Autonomous Dispatching Solution", "The closed-loop decision grid that runs predictive graphs and assists controllers")
    
    # 3 Solution Cards
    # Card 1: Telemetry
    draw_card(slide_3, Inches(0.8), Inches(2.0), Inches(3.64), Inches(4.5), COLOR_CARD, COLOR_CYAN, border_width_pt=1)
    t1 = slide_3.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(3.24), Inches(0.6))
    p = t1.text_frame.paragraphs[0]
    p.text = "TELEMETRY INGESTION"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    d1 = slide_3.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(3.24), Inches(3.5))
    tf = d1.text_frame
    tf.word_wrap = True
    tf.clear()
    add_bullet(tf, "GPS Sync: ", "Dynamically tracks active train spatial coordinates and speed profiles.")
    add_bullet(tf, "Block Status: ", "Monitors section block occupancy states (green, orange, red) in real-time.")
    add_bullet(tf, "Relational Store: ", "aiosqlite database writes state histories to ensure audit trail continuity.")

    # Card 2: Cascade Predictor
    draw_card(slide_3, Inches(4.84), Inches(2.0), Inches(3.64), Inches(4.5), COLOR_CARD, COLOR_CYAN, border_width_pt=1)
    t2 = slide_3.shapes.add_textbox(Inches(5.04), Inches(2.2), Inches(3.24), Inches(0.6))
    p = t2.text_frame.paragraphs[0]
    p.text = "CASCADE PREDICTIONS"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    d2 = slide_3.shapes.add_textbox(Inches(5.04), Inches(2.8), Inches(3.24), Inches(3.5))
    tf = d2.text_frame
    tf.word_wrap = True
    tf.clear()
    add_bullet(tf, "NetworkX Graphs: ", "Corridor stations and blocks represented as dynamic nodes and edges.")
    add_bullet(tf, "BFS Propagation: ", "Performs Breadth-First Searches to model delay transference down the corridor.")
    add_bullet(tf, "Conflict Warning: ", "Flags future route overlap conflicts up to 60 minutes in advance.")

    # Card 3: Explainable Resolvers
    draw_card(slide_3, Inches(8.88), Inches(2.0), Inches(3.64), Inches(4.5), COLOR_CARD, COLOR_CYAN, border_width_pt=1)
    t3 = slide_3.shapes.add_textbox(Inches(9.08), Inches(2.2), Inches(3.24), Inches(0.6))
    p = t3.text_frame.paragraphs[0]
    p.text = "EXPLAINABLE RESOLUTION"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    d3 = slide_3.shapes.add_textbox(Inches(9.08), Inches(2.8), Inches(3.24), Inches(3.5))
    tf = d3.text_frame
    tf.word_wrap = True
    tf.clear()
    add_bullet(tf, "Optimization: ", "Generates optimal loop-holding options to clear high-priority express routes.")
    add_bullet(tf, "Confidence Gate: ", "Automated actions execute if confidence > 85%, else escalates to operator.")
    add_bullet(tf, "Passenger RAC: ", "Forecasts waitlist confirmation chances for passengers delayed in cascades.")

    # ==========================================
    # SLIDE 4: Agentic AI Graph State Machine
    # ==========================================
    slide_4 = create_slide(prs)
    add_header(slide_4, "Agentic AI Graph State Machine", "Sequential collaboration pipeline of specialized agents resolving corridor bottlenecks")
    
    # 6 Agent boxes in a flowchart (3 columns x 2 rows)
    coords = [
        # Row 1
        {"name": "Monitor Agent", "desc": "Tracks Live Telemetry & Flags Variance > 20m", "left": 0.8, "top": 2.2},
        {"name": "Conflict Detector", "desc": "Scans Physical Corridor block Overlaps", "left": 4.84, "top": 2.2},
        {"name": "Cascade Predictor", "desc": "Runs BFS Delay Propagation downstream", "left": 8.88, "top": 2.2},
        # Row 2
        {"name": "Dispatch Agent", "desc": "Calculates Optimal Loop-Holding Actions", "left": 0.8, "top": 4.7},
        {"name": "Notification Agent", "desc": "Runs RAC waitlist calculations", "left": 4.84, "top": 4.7},
        {"name": "Audit Agent", "desc": "Seals steps in SHA-256 Block Ledger", "left": 8.88, "top": 4.7}
    ]
    
    for i, c in enumerate(coords):
        # Card background
        draw_card(slide_4, Inches(c["left"]), Inches(c["top"]), Inches(3.64), Inches(1.8), COLOR_CARD, COLOR_CYAN, border_width_pt=1)
        
        # Text inside
        tb = slide_4.shapes.add_textbox(Inches(c["left"] + 0.15), Inches(c["top"] + 0.25), Inches(3.34), Inches(1.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = c["name"].upper()
        p.font.name = "Segoe UI"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = COLOR_CYAN
        
        p2 = tf.add_paragraph()
        p2.space_before = Pt(6)
        p2.text = c["desc"]
        p2.font.name = "Segoe UI"
        p2.font.size = Pt(12)
        p2.font.color.rgb = COLOR_WHITE
        
    # Flow Arrows (horizontal & vertical transition)
    # Monitor -> Conflict (Row 1)
    arrow1 = slide_4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(2.9), Inches(0.28), Inches(0.2))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_BLUE
    arrow1.line.fill.background()
    
    # Conflict -> Cascade (Row 1)
    arrow2 = slide_4.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.54), Inches(2.9), Inches(0.28), Inches(0.2))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_BLUE
    arrow2.line.fill.background()
    
    # Cascade -> Audit (Down)
    arrow3 = slide_4.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.7), Inches(4.1), Inches(0.2), Inches(0.5))
    arrow3.fill.solid()
    arrow3.fill.fore_color.rgb = COLOR_BLUE
    arrow3.line.fill.background()
    
    # Audit -> Notification (Row 2 Left)
    arrow4 = slide_4.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(8.54), Inches(5.5), Inches(0.28), Inches(0.2))
    arrow4.fill.solid()
    arrow4.fill.fore_color.rgb = COLOR_BLUE
    arrow4.line.fill.background()
    
    # Notification -> Dispatch (Row 2 Left)
    arrow5 = slide_4.shapes.add_shape(MSO_SHAPE.LEFT_ARROW, Inches(4.5), Inches(5.5), Inches(0.28), Inches(0.2))
    arrow5.fill.solid()
    arrow5.fill.fore_color.rgb = COLOR_BLUE
    arrow5.line.fill.background()

    # ==========================================
    # SLIDE 5: Demo Scenario (NDLS Exit Failure)
    # ==========================================
    slide_5 = create_slide(prs)
    add_header(slide_5, "Demo Scenario: NDLS Signal Failure", "Timeline of the 7-step telemetry conflict resolution sequence")
    
    # Step Cards
    steps = [
        {"num": "STEP 1-2", "title": "DISRUPTION INGESTION", "detail": "Signal failure at NDLS exit tracks Shatabdi 12002 halted for 35m. ConflictDetector registers path block overlap.", "left": 0.8, "color": COLOR_RED},
        {"num": "STEP 3-4", "title": "CASCADE PROPAGATION", "detail": "CascadePredictor runs BFS. Forecasts Coal Freight BOXN-902 blocking NDLS-GZB corridor, projecting a massive +180m delay cascade.", "left": 4.84, "color": COLOR_RED},
        {"num": "STEP 5-7", "title": "RESOLUTION & AUDIT", "detail": "DispatchAgent recommends loop holding. Confidence is 78% (< 85% gate). Operator manual override approved and signed in ledger.", "left": 8.88, "color": COLOR_GREEN}
    ]
    
    for s in steps:
        draw_card(slide_5, Inches(s["left"]), Inches(2.2), Inches(3.64), Inches(4.3), COLOR_CARD, s["color"], border_width_pt=1.5)
        
        # Step Number
        tb_num = slide_5.shapes.add_textbox(Inches(s["left"] + 0.2), Inches(2.4), Inches(3.24), Inches(0.5))
        p_num = tb_num.text_frame.paragraphs[0]
        p_num.text = s["num"]
        p_num.font.name = "Segoe UI"
        p_num.font.size = Pt(18)
        p_num.font.bold = True
        p_num.font.color.rgb = s["color"]
        
        # Title
        tb_title = slide_5.shapes.add_textbox(Inches(s["left"] + 0.2), Inches(2.9), Inches(3.24), Inches(0.5))
        p_title = tb_title.text_frame.paragraphs[0]
        p_title.text = s["title"]
        p_title.font.name = "Segoe UI"
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        
        # Details
        tb_detail = slide_5.shapes.add_textbox(Inches(s["left"] + 0.2), Inches(3.5), Inches(3.24), Inches(2.7))
        tf_detail = tb_detail.text_frame
        tf_detail.word_wrap = True
        p_detail = tf_detail.paragraphs[0]
        p_detail.text = s["detail"]
        p_detail.font.name = "Segoe UI"
        p_detail.font.size = Pt(13)
        p_detail.font.color.rgb = COLOR_MUTED
        
    # Flow Lines
    line1 = slide_5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(4.44), Inches(4.35), Inches(0.4), Inches(0.02))
    line1.fill.solid()
    line1.fill.fore_color.rgb = COLOR_BLUE
    line1.line.fill.background()
    
    line2 = slide_5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(8.48), Inches(4.35), Inches(0.4), Inches(0.02))
    line2.fill.solid()
    line2.fill.fore_color.rgb = COLOR_BLUE
    line2.line.fill.background()

    # ==========================================
    # SLIDE 6: Decoupled Production Architecture
    # ==========================================
    slide_6 = create_slide(prs)
    add_header(slide_6, "Decoupled Full-Stack Architecture", "Dual production deployment strategy designed for performance and security isolation")
    
    # Left Card: Vite React Console
    draw_card(slide_6, Inches(0.8), Inches(2.0), Inches(5.66), Inches(4.5), COLOR_CARD, COLOR_CYAN, border_width_pt=1)
    
    t_left = slide_6.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.26), Inches(0.5))
    p = t_left.text_frame.paragraphs[0]
    p.text = "FRONTEND DASHBOARD CONSOLE (VERCEL)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    d_left = slide_6.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(5.26), Inches(3.5))
    tf_left = d_left.text_frame
    tf_left.word_wrap = True
    tf_left.clear()
    add_bullet(tf_left, "Interactive Corridor Map: ", "SVG render with flashing delay indicators.")
    add_bullet(tf_left, "Agent Terminal Stream: ", "Real-time process logs with dynamic tag filters.")
    add_bullet(tf_left, "Waitlist Solver UI: ", "Waitlist verification interface with visual progress circles.")
    add_bullet(tf_left, "Zero-Trust Sweeps: ", "Visual matrix scan verifying sealed audit ledger blocks.")
    add_bullet(tf_left, "Proxy Mappings: ", "vercel.json handles API routes directly to bypass CORS blocks.")

    # Right Card: FastAPI Backend
    draw_card(slide_6, Inches(6.86), Inches(2.0), Inches(5.66), Inches(4.5), COLOR_CARD, COLOR_BLUE, border_width_pt=1)
    
    t_right = slide_6.shapes.add_textbox(Inches(7.06), Inches(2.2), Inches(5.26), Inches(0.5))
    p = t_right.text_frame.paragraphs[0]
    p.text = "BACKEND INTELLIGENCE ENGINE (HUGGING FACE)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_BLUE
    
    d_right = slide_6.shapes.add_textbox(Inches(7.06), Inches(2.8), Inches(5.26), Inches(3.5))
    tf_right = d_right.text_frame
    tf_right.word_wrap = True
    tf_right.clear()
    add_bullet(tf_right, "API Web Framework: ", "Python FastAPI serving asynchronous, high-frequency endpoints.")
    add_bullet(tf_right, "NetworkX Analytics: ", "Direct graph-based evaluation of route bottlenecks.")
    add_bullet(tf_right, "Data Consistency: ", "SQLite relational database using aiosqlite transactions.")
    add_bullet(tf_right, "Dockerized Space: ", "Custom Dockerfile exposing server instance on port 7860.")
    add_bullet(tf_right, "Secured Hashes: ", "Seals controller manual approvals with SHA-256 cryptography.")

    # ==========================================
    # SLIDE 7: Cryptographic Decision Ledger
    # ==========================================
    slide_7 = create_slide(prs)
    add_header(slide_7, "Zero-Trust Cryptographic Audit Ledger", "Blockchain-inspired sequential integrity verification protecting grid operations")
    
    # 3 Ledgers Blocks
    blocks = [
        {"num": "GENESIS BLOCK 001", "action": "NDLS DISRUPTION", "hash": "0000a89d3c5f21ea", "prev": "----------------", "left": 0.8},
        {"num": "LEDGER BLOCK 002", "action": "DISPATCH PROPOSE", "hash": "0000ef82c1b2f90d", "prev": "0000a89d3c5f21ea", "left": 4.84},
        {"num": "LEDGER BLOCK 003", "action": "OPERATOR APPROVE", "hash": "00009c2af0e5138a", "prev": "0000ef82c1b2f90d", "left": 8.88}
    ]
    
    for b in blocks:
        draw_card(slide_7, Inches(b["left"]), Inches(2.4), Inches(3.64), Inches(4.0), COLOR_CARD, COLOR_GREEN, border_width_pt=1.5)
        
        # Block Header
        tb_head = slide_7.shapes.add_textbox(Inches(b["left"] + 0.2), Inches(2.6), Inches(3.24), Inches(0.5))
        p = tb_head.text_frame.paragraphs[0]
        p.text = b["num"]
        p.font.name = "Segoe UI"
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = COLOR_GREEN
        
        # Content
        tb_cont = slide_7.shapes.add_textbox(Inches(b["left"] + 0.2), Inches(3.2), Inches(3.24), Inches(3.0))
        tf = tb_cont.text_frame
        tf.word_wrap = True
        tf.clear()
        add_bullet(tf, "Action: ", b["action"], font_size=12)
        add_bullet(tf, "Prev Hash: ", b["prev"], font_size=12, color=COLOR_MUTED)
        add_bullet(tf, "Block Hash: ", b["hash"], font_size=12, color=COLOR_WHITE)
        add_bullet(tf, "Signature: ", "ECDSA Secured", font_size=12, color=COLOR_MUTED)
        
    # Block Links (Arrow indicators)
    arrow1 = slide_7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(4.5), Inches(4.2), Inches(0.28), Inches(0.2))
    arrow1.fill.solid()
    arrow1.fill.fore_color.rgb = COLOR_GREEN
    arrow1.line.fill.background()
    
    arrow2 = slide_7.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(8.54), Inches(4.2), Inches(0.28), Inches(0.2))
    arrow2.fill.solid()
    arrow2.fill.fore_color.rgb = COLOR_GREEN
    arrow2.line.fill.background()

    # ==========================================
    # SLIDE 8: Tokyo Grand Finals Roadmap
    # ==========================================
    slide_8 = create_slide(prs)
    add_header(slide_8, "Tokyo Grand Finals Roadmap", "Scaling RailMind from a regional pilot to a unified national signaling intelligence")
    
    # Left Card: Phase 1
    draw_card(slide_8, Inches(0.8), Inches(2.0), Inches(5.66), Inches(4.5), COLOR_CARD, COLOR_MUTED, border_width_pt=1)
    
    t_p1 = slide_8.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(5.26), Inches(0.5))
    p = t_p1.text_frame.paragraphs[0]
    p.text = "PHASE 1: DELHI REGIONALS (CURRENT)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_MUTED
    
    d_p1 = slide_8.shapes.add_textbox(Inches(1.0), Inches(2.8), Inches(5.26), Inches(3.5))
    tf_p1 = d_p1.text_frame
    tf_p1.word_wrap = True
    tf_p1.clear()
    add_bullet(tf_p1, "Seeded Database: ", "Relational mapping for the Delhi-Kanpur corridor.")
    add_bullet(tf_p1, "Interactive Web Console: ", "React client displaying live dispatch recommendations.")
    add_bullet(tf_p1, "Predictive Simulations: ", "Offline 7-step player mapping delay chains.")
    add_bullet(tf_p1, "Ledger Sweeps: ", "Local SHA-256 validation routines.")

    # Right Card: Phase 2
    draw_card(slide_8, Inches(6.86), Inches(2.0), Inches(5.66), Inches(4.5), COLOR_CARD, COLOR_CYAN, border_width_pt=1.5)
    
    t_p2 = slide_8.shapes.add_textbox(Inches(7.06), Inches(2.2), Inches(5.26), Inches(0.5))
    p = t_p2.text_frame.paragraphs[0]
    p.text = "PHASE 2: TOKYO GRAND FINALS (SCALE)"
    p.font.name = "Segoe UI"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = COLOR_CYAN
    
    d_p2 = slide_8.shapes.add_textbox(Inches(7.06), Inches(2.8), Inches(5.26), Inches(3.5))
    tf_p2 = d_p2.text_frame
    tf_p2.word_wrap = True
    tf_p2.clear()
    add_bullet(tf_p2, "Real-Time Ingestion: ", "Integrating Kafka streams for live signaling tracks.")
    add_bullet(tf_p2, "Partitioned Storage: ", "TimescaleDB telemetry engines handling millions of daily updates.")
    add_bullet(tf_p2, "XGBoost Ticket Models: ", "ML Waitlist probability inference engine trained on seasonal records.")
    add_bullet(tf_p2, "Kavach Protocol: ", "Direct API endpoints interacting with locomotive warning networks.")
    add_bullet(tf_p2, "Geographic Clusters: ", "High-availability Kubernetes deployment covering all primary zones.")

    # Save presentation
    script_dir = os.path.dirname(os.path.abspath(__file__))
    output_filename = os.path.join(script_dir, "../../railmind_presentation.pptx")
    prs.save(output_filename)
    print(f"Presentation successfully created and saved to: {output_filename}")

if __name__ == "__main__":
    build_deck()
