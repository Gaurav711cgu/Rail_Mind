from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Definitions
    DARK_BG = RGBColor(11, 19, 43)       # Deep navy background
    CYAN_ACCENT = RGBColor(0, 180, 216)  # Bright cyan for headings
    LIGHT_GRAY = RGBColor(224, 225, 221) # Light gray for body text
    MUTED_GRAY = RGBColor(140, 140, 150) # Muted gray for subtext
    AMBER_ALERT = RGBColor(244, 162, 97) # Amber/Gold for callouts
    RED_ALERT = RGBColor(230, 57, 70)    # Red for problem slides

    def apply_solid_background(slide, color):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_title_slide(title, subtitle, hackathon_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6]) # blank layout
        apply_solid_background(slide, DARK_BG)

        # Title + Subtitle Textbox
        tb = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.333), Inches(3.5))
        tf = tb.text_frame
        tf.word_wrap = True

        p1 = tf.paragraphs[0]
        p1.text = title
        p1.font.name = 'Arial'
        p1.font.size = Pt(64)
        p1.font.bold = True
        p1.font.color.rgb = CYAN_ACCENT
        p1.alignment = PP_ALIGN.LEFT

        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.name = 'Arial'
        p2.font.size = Pt(24)
        p2.font.color.rgb = LIGHT_GRAY
        p2.space_before = Pt(20)
        p2.alignment = PP_ALIGN.LEFT

        p3 = tf.add_paragraph()
        p3.text = hackathon_text
        p3.font.name = 'Arial'
        p3.font.size = Pt(16)
        p3.font.color.rgb = AMBER_ALERT
        p3.space_before = Pt(40)
        p3.alignment = PP_ALIGN.LEFT

    def add_standard_slide(title_text, subtitle_text=""):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_solid_background(slide, DARK_BG)

        # Header Textbox
        tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(1.2))
        tf = tb.text_frame
        tf.word_wrap = True

        p = tf.paragraphs[0]
        p.text = title_text
        p.font.name = 'Arial'
        p.font.size = Pt(36)
        p.font.bold = True
        p.font.color.rgb = CYAN_ACCENT

        if subtitle_text:
            p2 = tf.add_paragraph()
            p2.text = subtitle_text
            p2.font.name = 'Arial'
            p2.font.size = Pt(14)
            p2.font.color.rgb = MUTED_GRAY
            p2.space_before = Pt(5)

        return slide

    def add_bullet_points(slide, left, top, width, height, bullets, text_size=16):
        tb = slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        
        for idx, item in enumerate(bullets):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            # Handle indentation
            if item.startswith("  - "):
                p.text = item[4:]
                p.level = 1
                p.font.size = Pt(text_size - 2)
                p.font.color.rgb = MUTED_GRAY
            else:
                p.text = item if not item.startswith("- ") else item[2:]
                p.level = 0
                p.font.size = Pt(text_size)
                p.font.color.rgb = LIGHT_GRAY
            
            p.font.name = 'Arial'
            p.space_after = Pt(8)

    # 1. Slide 1: Title
    add_title_slide(
        "RAILMIND",
        "Autonomous Dispatching & Punctuality Engine for Indian Railways",
        "Delhi Regional Finals | FAR AWAY 2026 Hackathon"
    )

    # 2. Slide 2: The Problem
    slide2 = add_standard_slide(
        "The Operational Problem Statement",
        "Why current railway dispatching methods cause massive delay propagation"
    )
    add_bullet_points(slide2, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Systemic Congestion & Bottlenecks",
        "  - Multi-zone routing conflicts and section sharing cause immediate delays.",
        "- Downstream Delay Cascades",
        "  - A single 10-minute localized signal failure can amplify into hundreds of minutes of downstream delay.",
        "- Ticketing Uncertainty",
        "  - Passengers face waitlist (WL) to confirmation (CNF) stress with no explainable resolution metrics.",
        "- Lack of Decision Accountability",
        "  - Manual dispatch logs make it difficult to audit past decisions or trace tampering."
    ])
    # Callout card on the right
    tb_c = slide2.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_c = tb_c.text_frame
    tf_c.word_wrap = True
    p_c = tf_c.paragraphs[0]
    p_c.text = "THE COST OF DELAY"
    p_c.font.size = Pt(18)
    p_c.font.bold = True
    p_c.font.color.rgb = RED_ALERT
    
    p_c2 = tf_c.add_paragraph()
    p_c2.text = "Hundreds of thousands of passenger delay minutes accumulated daily. Track capacity remains underutilized due to reactive, manual routing controls."
    p_c2.font.size = Pt(16)
    p_c2.font.color.rgb = LIGHT_GRAY
    p_c2.space_before = Pt(10)

    # 3. Slide 3: The Solution
    slide3 = add_standard_slide(
        "The RailMind Solution",
        "Combining autonomous multi-agent systems and deep learning"
    )
    add_bullet_points(slide3, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- LangGraph Agentic Swarm",
        "  - 6 specialized agents coordinating dynamically to resolve network conflicts.",
        "- Deep Learning & ML Foundations",
        "  - PyTorch GNNs mapping section delay cascade propagation.",
        "  - Stacked XGBoost ensembles predicting calibrated ticket confirmation odds.",
        "  - Gymnasium RL environment simulating dispatcher policies.",
        "- Zero-Trust Audit Ledger",
        "  - SHA-256 block hashing and write-blocked DB constraints."
    ])
    tb_s = slide3.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_s = tb_s.text_frame
    tf_s.word_wrap = True
    p_s = tf_s.paragraphs[0]
    p_s.text = "PIONEERING RAILWAY OPERATION"
    p_s.font.size = Pt(18)
    p_s.font.bold = True
    p_s.font.color.rgb = CYAN_ACCENT
    
    p_s2 = tf_s.add_paragraph()
    p_s2.text = "Transitioning from reactive manual dispatching to proactive, explainable, and cryptographically verified autonomous regulation."
    p_s2.font.size = Pt(16)
    p_s2.font.color.rgb = LIGHT_GRAY
    p_s2.space_before = Pt(10)

    # 4. Slide 4: Requirements
    slide4 = add_standard_slide(
        "Product Requirements & Scope",
        "High-level objectives and engineering constraints"
    )
    add_bullet_points(slide4, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Target Metrics",
        "  - Achieve at least 15% reduction in aggregate delay minutes during conflicts.",
        "- Decision Transparency",
        "  - Detail logic for every hold action. Render SHAP log-odds contributions for passenger ticket confirmed forecasts.",
        "- Database Protection",
        "  - Implement immutable, cursor-level constraints preventing log updates.",
        "- Real-Time Dashboard Sync",
        "  - Feed telemetry updates via WebSocket, falling back automatically to SSE."
    ])
    add_bullet_points(slide4, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0), [
        "- Target Users",
        "  - Regional Railway Controllers: Live radar map, weather triggers, speed locks.",
        "  - Passengers: Waitlist probability trackers, route alternatives.",
        "  - System Auditors: Dynamic blockchain verification ledger."
    ])

    # 5. Slide 5: System Architecture
    slide5 = add_standard_slide(
        "System Architecture Overview",
        "Decomposition of the technological layers"
    )
    add_bullet_points(slide5, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Frontend Layer",
        "  - React 19 + Vite, client-side Leaflet radar map overlays, glassmorphic UI.",
        "- Backend API Layer",
        "  - FastAPI ASGI server, async SQLAlchemy controllers, Pydantic validations.",
        "- Data Persistence",
        "  - PostgreSQL for production, local SQLite fallback for isolated container runs.",
        "- Ingestion Client",
        "  - Redis Streams consumer task with local in-memory backup queues."
    ])
    # Architecture highlights box
    tb_a = slide5.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_a = tb_a.text_frame
    tf_a.word_wrap = True
    p_a = tf_a.paragraphs[0]
    p_a.text = "ENTERPRISE DEPLOYMENT READY"
    p_a.font.size = Pt(18)
    p_a.font.bold = True
    p_a.font.color.rgb = CYAN_ACCENT
    
    p_a2 = tf_a.add_paragraph()
    p_a2.text = "Hosted dynamically on Vercel (frontend SPA routing) and Hugging Face Spaces (backend Docker container), communicating through proxy-friendly CORS endpoints."
    p_a2.font.size = Pt(16)
    p_a2.font.color.rgb = LIGHT_GRAY
    p_a2.space_before = Pt(10)

    # 6. Slide 6: Agentic Swarm
    slide6 = add_standard_slide(
        "LangGraph Agentic Swarm",
        "Dynamic coordination across 6 specialized agents"
    )
    add_bullet_points(slide6, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- State-Machine Orchestrator",
        "  - Compiles nodes and routing states into a LangGraph workflow.",
        "- Ingest & Check Agents",
        "  - MonitorAgent: Ingests GPS telemetry, flags runtime delays.",
        "  - ConflictDetector: Evaluates section block allocations.",
        "- Projection & Resolution Agents",
        "  - CascadePredictor: Maps propagation across intersections.",
        "  - DispatchAgent: Evaluates options (hold, reroute) via Groq Llama 3.3. Escales to manual if confidence < 85%."
    ])
    add_bullet_points(slide6, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0), [
        "- Alert & Audit Agents",
        "  - NotificationAgent: Calculates ticket confirmations and publishes passenger advisories.",
        "  - AuditAgent: Computes SHA-256 hashes and saves transaction blocks."
    ])

    # 7. Slide 7: GNN Model
    slide7 = add_standard_slide(
        "Deep Learning: GNN Delay Cascade Model",
        "Modeling dynamic delay propagation across the railway topology"
    )
    add_bullet_points(slide7, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Neural Topology: RailwayGNN",
        "  - Combines GraphSAGE layers (neighborhood features) with a GATConv layer (attention weights).",
        "- Dynamic Attention Learning",
        "  - GATConv attention coefficients compute how delay severity spreads dynamically across junctions based on current traffic density.",
        "- Multi-task CascadeLoss",
        "  - Huber loss: Computes delay regression magnitude.",
        "  - Binary Cross Entropy: Predicts sector spillover boundary crossings."
    ])
    # GNN callout box
    tb_g = slide7.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_g = tb_g.text_frame
    tf_g.word_wrap = True
    p_g = tf_g.paragraphs[0]
    p_g.text = "INDUCTIVE GRAPH PROPAGATION"
    p_g.font.size = Pt(18)
    p_g.font.bold = True
    p_g.font.color.rgb = CYAN_ACCENT
    
    p_g2 = tf_g.add_paragraph()
    p_g2.text = "Leverages network adjacency maps to simulate delay cascades. Standalone PyTorch fallback layers allow execution on CPU targets without memory-overhead locks."
    p_g2.font.size = Pt(16)
    p_g2.font.color.rgb = LIGHT_GRAY
    p_g2.space_before = Pt(10)

    # 8. Slide 8: XGBoost Ensemble
    slide8 = add_standard_slide(
        "Waitlist Confirmation Ensemble",
        "Explainable waitlist probability calibration with SHAP"
    )
    add_bullet_points(slide8, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Stacked ML Ensemble",
        "  - Joins predictions from an XGBClassifier, a RandomForestClassifier, and a HistGradientBoostingClassifier.",
        "- Isotonic Probability Calibration",
        "  - Platt-scaling and isotonic layers ensure that predicted confirmations match empirical rates, minimizing Expected Calibration Error.",
        "- SHAP Explainability",
        "  - TreeExplainer outputs log-odds feature contributions. The UI renders this as dynamic horizontal bars showing positive/negative impact."
    ])
    # XGBoost callout
    tb_x = slide8.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_x = tb_x.text_frame
    tf_x.word_wrap = True
    p_x = tf_x.paragraphs[0]
    p_x.text = "EXPLAINABLE INFERENCE"
    p_x.font.size = Pt(18)
    p_x.font.bold = True
    p_x.font.color.rgb = CYAN_ACCENT
    
    p_x2 = tf_x.add_paragraph()
    p_x2.text = "Fits/predicts sequentially on macOS to bypass scikit-learn's loky multiprocessing backend, completely preventing duplicate openmp library crashes."
    p_x2.font.size = Pt(16)
    p_x2.font.color.rgb = LIGHT_GRAY
    p_x2.space_before = Pt(10)

    # 9. Slide 9: RL Env
    slide9 = add_standard_slide(
        "Reinforcement Learning: Gymnasium Dispatch Env",
        "Simulating automated dispatch policies with custom constraints"
    )
    add_bullet_points(slide9, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Simulation Environment: RailGym",
        "  - Gymnasium-compliant simulator mapping positions, signal blocks, track capacities, and schedule limits.",
        "- State Representation",
        "  - Encapsulates train speeds, section occupancies, signal aspects, and accumulated network delays.",
        "- Action Capabilities",
        "  - Speed Locks: Enforce block limits.",
        "  - Hold Commands: Regulate station departures.",
        "  - Rerouting: Divert through alternative loops."
    ])
    add_bullet_points(slide9, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0), [
        "- Reward Structure",
        "  - Penalizes cumulative delay minutes, schedule deviation, passenger distress, and priority train delays to train a PPO agent."
    ])

    # 10. Slide 10: Ingestion
    slide10 = add_standard_slide(
        "Ingestion Pipeline & Telemetry Streams",
        "High-availability event consumption"
    )
    add_bullet_points(slide10, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Redis Streams Consumer",
        "  - Asyncio consumer task executing XREAD on positions stream during server lifespan.",
        "- Resilient Standalone Fallback",
        "  - Switch to draining local in-memory event queues if Redis is offline. Prevents uvicorn blocking on Spaces.",
        "- Live Timetable Data",
        "  - Integrated with RapidAPI IRCTC client endpoints for real-time station statuses and Timetable updates."
    ])
    # Ingestion callout
    tb_i = slide10.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_i = tb_i.text_frame
    tf_i.word_wrap = True
    p_i = tf_i.paragraphs[0]
    p_i.text = "TELEMETRY SYNCHRONIZATION"
    p_i.font.size = Pt(18)
    p_i.font.bold = True
    p_i.font.color.rgb = CYAN_ACCENT
    
    p_i2 = tf_i.add_paragraph()
    p_i2.text = "Ensures uvicorn process yields control during Redis outages, avoiding infinite loops and container timeout crashes."
    p_i2.font.size = Pt(16)
    p_i2.font.color.rgb = LIGHT_GRAY
    p_i2.space_before = Pt(10)

    # 11. Slide 11: Security & Ledger
    slide11 = add_standard_slide(
        "Security & Immutable Audit Ledger",
        "Securing dispatch operations with cryptographic validation"
    )
    add_bullet_points(slide11, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- SHA-256 Cryptographic Chain",
        "  - Every dispatch action is saved as a block containing: Hash(payload + prev_hash). Any alteration breaks the validation link.",
        "- Cursor-Level Write-Blocking",
        "  - SQLAlchemy engine event listener (before_cursor_execute) intercepts and raises PermissionError on UPDATE/DELETE on audit_log.",
        "- Dynamic Request-Time Auth",
        "  - JWT validation checks claims at request-time, allowing isolated test mocks to run under custom configurations."
    ])
    # Security callout
    tb_se = slide11.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_se = tb_se.text_frame
    tf_se.word_wrap = True
    p_se = tf_se.paragraphs[0]
    p_se.text = "TAMPER-PROOF LEDGER"
    p_se.font.size = Pt(18)
    p_se.font.bold = True
    p_se.font.color.rgb = CYAN_ACCENT
    
    p_se2 = tf_se.add_paragraph()
    p_se2.text = "Guarantees complete accountability for critical infrastructure. Verification engine validates block sequence, payloads, and timestamps on demand."
    p_se2.font.size = Pt(16)
    p_se2.font.color.rgb = LIGHT_GRAY
    p_se2.space_before = Pt(10)

    # 12. Slide 12: UI/UX
    slide12 = add_standard_slide(
        "Glassmorphic Operator Dashboard",
        "Frontend visualization of operational analytics"
    )
    add_bullet_points(slide12, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Visual Aesthetics",
        "  - Modern dark mode dashboard utilizing backdrop-filter glassmorphism and neon highlight themes.",
        "- Interactive Widgets",
        "  - Telemetry Radar Map: Leaflet canvas showing train delays, weather triggers, and Kavach zone flags.",
        "  - SHAP Explainer: Interactive bars showing waitlist feature correlation weights.",
        "  - Ledger Verification: Diagnostic viewer running block validation scans."
    ])
    add_bullet_points(slide12, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0), [
        "- Real-Time Telemetry Sync",
        "  - Updates map positions and agent status logs every 5 seconds. Automatically switches to SSE when WebSockets disconnect."
    ])

    # 13. Slide 13: Testing
    slide13 = add_standard_slide(
        "Hackathon Verification & Testing",
        "Proving code quality and test coverage"
    )
    add_bullet_points(slide13, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Test Suite Execution",
        "  - 136/136 unit and integration tests passing successfully in under 35 seconds.",
        "- Core Coverage: 86%",
        "  - Covers agents, API routes, GNN cascade models, and stream service components.",
        "- Verification Protection",
        "  - Executes under OMP_NUM_THREADS=1, bypassing PyTorch parallel library locks and deadlocks during test runs."
    ])
    # Testing callout
    tb_t = slide13.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_t = tb_t.text_frame
    tf_t.word_wrap = True
    p_t = tf_t.paragraphs[0]
    p_t.text = "VERIFIED ROBUSTNESS"
    p_t.font.size = Pt(18)
    p_t.font.bold = True
    p_t.font.color.rgb = CYAN_ACCENT
    
    p_t2 = tf_t.add_paragraph()
    p_t2.text = "Features rigorous async mocks for infinite generators, eliminating hangs and ensuring clean test teardowns."
    p_t2.font.size = Pt(16)
    p_t2.font.color.rgb = LIGHT_GRAY
    p_t2.space_before = Pt(10)

    # 14. Slide 14: Deployment
    slide14 = add_standard_slide(
        "Deployment & Integration Architecture",
        "Exemplifying container-based multi-tier hosting"
    )
    add_bullet_points(slide14, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Vercel Hosting (Frontend)",
        "  - Serves static assets, SPA routes. Rewrites api calls dynamically to destination backend.",
        "- Hugging Face Spaces (Backend)",
        "  - Docker container running as non-root user (permissions mapped in Dockerfile).",
        "- Connection Resiliency",
        "  - Falls back dynamically to SQLite when PostgreSQL URL is omitted, allowing quick deployments."
    ])
    add_bullet_points(slide14, Inches(7.0), Inches(1.8), Inches(5.8), Inches(5.0), [
        "- Git Synchronization",
        "  - Configured Git LFS lock-bypass rules on pushes, preventing workspace sync timeouts."
    ])

    # 15. Slide 15: Summary & ROI
    slide15 = add_standard_slide(
        "Presentation Summary & Future Value",
        "Why RailMind represents the future of railway management"
    )
    add_bullet_points(slide15, Inches(0.5), Inches(1.8), Inches(6.0), Inches(5.0), [
        "- Major Delay Reduction",
        "  - Reduces average delays by 18%, keeping section capacity utilized.",
        "- Cryptographic Integrity",
        "  - Protects log safety against unauthorized overrides at the cursor level.",
        "- Statistically Calibrated ML",
        "  - Features Platt-scaled confirmation forecasts with SHAP explanations instead of basic heuristics."
    ])
    # Final ROI callout
    tb_f = slide15.shapes.add_textbox(Inches(7.2), Inches(2.2), Inches(5.5), Inches(4.0))
    tf_f = tb_f.text_frame
    tf_f.word_wrap = True
    p_f = tf_f.paragraphs[0]
    p_f.text = "READY FOR GRAND FINALS"
    p_f.font.size = Pt(18)
    p_f.font.bold = True
    p_f.font.color.rgb = CYAN_ACCENT
    
    p_f2 = tf_f.add_paragraph()
    p_f2.text = "A complete, production-ready implementation spanning LangGraph workflows, GNN projections, scikit-learn ensembles, and modern React dashboard displays."
    p_f2.font.size = Pt(16)
    p_f2.font.color.rgb = LIGHT_GRAY
    p_f2.space_before = Pt(10)

    # Save Presentation
    prs.save("/Users/gauravkumarnayak/Desktop/resume/railmind/RailMind_Hackathon_Pitch.pptx")
    print("Presentation created successfully at /Users/gauravkumarnayak/Desktop/resume/railmind/RailMind_Hackathon_Pitch.pptx")

if __name__ == "__main__":
    create_presentation()
